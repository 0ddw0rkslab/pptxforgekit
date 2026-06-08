from __future__ import annotations

import logging
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from pptxforgekit.chart.builder import ChartBuilder
from pptxforgekit.exceptions import RenderError
from pptxforgekit.models.schema import (
    ChartElement,
    ImageElement,
    PresentationSchema,
    SlideSchema,
    TableElement,
    TextElement,
    TextStyle,
)
from pptxforgekit.models.theme import ThemeConfig
from pptxforgekit.renderer.layout_engine import get_positions

logger = logging.getLogger(__name__)

_SLIDE_WIDTH_CM = 33.87
_SLIDE_HEIGHT_CM = 19.05

_ALIGN_MAP = {
    "left":    PP_ALIGN.LEFT,
    "center":  PP_ALIGN.CENTER,
    "right":   PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _sorted_by_z(*element_lists: list) -> list:
    """Merge element lists and sort by z_index ascending."""
    merged = [e for lst in element_lists for e in lst]
    return sorted(merged, key=lambda e: getattr(e, "z_index", 0))


class PPTXRenderer:
    def __init__(self, theme: ThemeConfig) -> None:
        self._theme = theme
        self._chart_builder = ChartBuilder()

    def render(
        self,
        schema: PresentationSchema,
        output_path: Path,
        base_path: Path | None = None,
    ) -> None:
        logger.info("Rendering %d slides → %s", len(schema.slides), output_path)
        prs = self._new_presentation()
        blank_layout = prs.slide_layouts[6]  # "Blank" layout index

        for slide_schema in schema.slides:
            slide = prs.slides.add_slide(blank_layout)
            resolved_base = base_path or output_path.parent
            try:
                self._render_slide(slide, slide_schema, resolved_base)
            except Exception as exc:
                raise RenderError(
                    f"Failed to render slide '{slide_schema.slide_id}': {exc}"
                ) from exc

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info("Saved: %s", output_path)

    # ── slide-level ────────────────────────────────────────────────────────────

    def _new_presentation(self) -> Presentation:
        prs = Presentation()
        prs.slide_width = Cm(_SLIDE_WIDTH_CM)
        prs.slide_height = Cm(_SLIDE_HEIGHT_CM)
        return prs

    def _normalize_overlaps(self, schema: SlideSchema) -> SlideSchema:
        """Redistribute overlapping text elements before rendering.

        When multiple non-title text elements share the same position (a common
        LLM generation artefact), they are either redistributed into the
        layout's column slots (two_column) or merged into a single element.
        Title elements are always passed through unchanged.
        """
        pos_groups: defaultdict[tuple[float, float, float, float], list[TextElement]] = defaultdict(list)
        title_elems: list[TextElement] = []

        for elem in schema.text_elements:
            if elem.role == "title":
                title_elems.append(elem)
                continue
            p = elem.position
            key = (round(p.x, 1), round(p.y, 1), round(p.w, 1), round(p.h, 1))
            pos_groups[key].append(elem)

        # Fast path: no overlaps
        if all(len(g) == 1 for g in pos_groups.values()):
            return schema

        layout_pos = get_positions(schema.layout_type)
        new_elems: list[TextElement] = list(title_elems)

        for elems in pos_groups.values():
            if len(elems) == 1:
                new_elems.append(elems[0])
            elif "body_left" in layout_pos and "body_right" in layout_pos:
                # two_column: first block → left column, remaining → right column
                new_elems.append(elems[0].model_copy(
                    update={"position": layout_pos["body_left"]}
                ))
                right_content = "\n".join(e.content for e in elems[1:])
                new_elems.append(elems[1].model_copy(
                    update={"content": right_content, "position": layout_pos["body_right"]}
                ))
            else:
                # Other layouts: merge all content into the first element
                merged = "\n".join(e.content for e in elems)
                new_elems.append(elems[0].model_copy(update={"content": merged}))

        logger.debug(
            "Slide %s: normalized %d overlapping text elements",
            schema.slide_id,
            sum(len(g) for g in pos_groups.values() if len(g) > 1),
        )
        return schema.model_copy(update={"text_elements": new_elems})

    def _render_slide(
        self, slide: object, schema: SlideSchema, base_path: Path
    ) -> None:
        schema = self._normalize_overlaps(schema)
        bg = schema.background_color or self._theme.colors.background
        self._set_background(slide, bg)

        # Render elements in z_index order
        for elem in _sorted_by_z(
            schema.text_elements,
            schema.image_elements,
            schema.chart_elements,
            schema.table_elements,
        ):
            if isinstance(elem, TextElement):
                self._add_text(slide, elem)               # type: ignore[arg-type]
            elif isinstance(elem, ChartElement):
                self._add_chart(slide, elem, base_path)   # type: ignore[arg-type]
            elif isinstance(elem, ImageElement):
                self._add_image(slide, elem, base_path)   # type: ignore[arg-type]
            elif isinstance(elem, TableElement):
                self._add_table(slide, elem)              # type: ignore[arg-type]

        if schema.speaker_note:
            self._add_speaker_note(slide, schema.speaker_note)  # type: ignore[arg-type]

        self._add_footer(slide, schema.footer_override)         # type: ignore[arg-type]

    def _set_background(self, slide: object, hex_color: str) -> None:
        fill = slide.background.fill  # type: ignore[union-attr]
        fill.solid()
        fill.fore_color.rgb = _rgb(hex_color)

    # ── text ──────────────────────────────────────────────────────────────────

    def _add_text(self, slide: object, elem: TextElement) -> None:
        p = elem.position
        tx_box = slide.shapes.add_textbox(  # type: ignore[union-attr]
            Cm(p.x), Cm(p.y), Cm(p.w), Cm(p.h)
        )
        tf = tx_box.text_frame
        tf.word_wrap = True

        style = elem.style
        lines = elem.content.splitlines()
        first_para = True
        for line in lines:
            if first_para:
                para = tf.paragraphs[0]
                first_para = False
            else:
                para = tf.add_paragraph()
            para.clear()
            run = para.add_run()
            run.text = line
            self._apply_run_style(run, style, elem.role)
            para.alignment = _ALIGN_MAP.get(style.align, PP_ALIGN.LEFT)

    def _apply_run_style(self, run: object, style: TextStyle, role: str) -> None:
        run.font.size = Pt(style.font_size)          # type: ignore[union-attr]
        run.font.bold = style.bold                   # type: ignore[union-attr]
        run.font.italic = style.italic               # type: ignore[union-attr]
        run.font.underline = style.underline         # type: ignore[union-attr]
        color = style.color or self._theme.colors.text
        run.font.color.rgb = _rgb(color)             # type: ignore[union-attr]
        run.font.name = (                            # type: ignore[union-attr]
            self._theme.fonts.title_font
            if role in ("title", "subtitle")
            else self._theme.fonts.body_font
        )

    # ── chart ─────────────────────────────────────────────────────────────────

    def _add_chart(self, slide: object, elem: ChartElement, base_path: Path) -> None:
        try:
            self._chart_builder.add_chart_to_slide(
                slide,   # type: ignore[arg-type]
                elem,
                self._theme,
                base_path=base_path,
            )
        except Exception as exc:
            logger.warning(
                "Chart '%s' could not be rendered (%s). Inserting placeholder.",
                elem.element_id, exc,
            )
            self._add_placeholder(slide, elem.position, f"[Chart: {elem.chart_type} — {elem.data_source}]")  # type: ignore[arg-type]

    # ── image ─────────────────────────────────────────────────────────────────

    def _add_image(self, slide: object, elem: ImageElement, base_path: Path) -> None:
        img_path = Path(elem.file_path)
        if not img_path.is_absolute():
            img_path = base_path / img_path
        if not img_path.exists():
            logger.warning("Image not found, skipping: %s", img_path)
            self._add_placeholder(
                slide, elem.position,   # type: ignore[arg-type]
                f"[Image: {elem.file_path}]",
            )
            return

        p = elem.position
        slide.shapes.add_picture(str(img_path), Cm(p.x), Cm(p.y), Cm(p.w), Cm(p.h))  # type: ignore[union-attr]

        if elem.figure_label:
            self._add_figure_label(slide, elem)  # type: ignore[arg-type]

        if elem.caption:
            self._add_caption(slide, elem)       # type: ignore[arg-type]

    def _add_figure_label(self, slide: object, elem: ImageElement) -> None:
        p = elem.position
        lbl = slide.shapes.add_textbox(  # type: ignore[union-attr]
            Cm(p.x), Cm(p.y), Cm(2.0), Cm(0.8)
        )
        tf = lbl.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = elem.figure_label
        run.font.size = Pt(self._theme.fonts.caption_size)
        run.font.bold = True
        run.font.color.rgb = _rgb(self._theme.colors.primary)

    def _add_caption(self, slide: object, elem: ImageElement) -> None:
        p = elem.position
        cap = slide.shapes.add_textbox(  # type: ignore[union-attr]
            Cm(p.x), Cm(p.y + p.h + 0.1), Cm(p.w), Cm(0.9)
        )
        tf = cap.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = elem.caption
        run.font.size = Pt(self._theme.fonts.caption_size)
        run.font.italic = True
        run.font.color.rgb = _rgb(self._theme.colors.text_light)

    # ── table ─────────────────────────────────────────────────────────────────

    def _add_table(self, slide: object, elem: TableElement) -> None:
        n_cols = len(elem.headers)
        if n_cols == 0:
            return
        n_rows = len(elem.rows) + 1   # +1 for header row

        p = elem.position
        tbl_frame = slide.shapes.add_table(  # type: ignore[union-attr]
            n_rows, n_cols, Cm(p.x), Cm(p.y), Cm(p.w), Cm(p.h)
        )
        tbl = tbl_frame.table

        # Apply per-column widths if specified
        if elem.column_widths:
            for j, cw in enumerate(elem.column_widths):
                tbl.columns[j].width = Cm(cw)

        header_bg = _rgb(self._theme.colors.primary)
        header_fg = _rgb("#FFFFFF")
        body_size = Pt(self._theme.fonts.body_size - 2)

        for j, header in enumerate(elem.headers):
            cell = tbl.cell(0, j)
            para = cell.text_frame.paragraphs[0]
            para.clear()
            run = para.add_run()
            run.text = str(header)
            run.font.bold = True
            run.font.color.rgb = header_fg
            run.font.size = body_size
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg

        for i, row_data in enumerate(elem.rows):
            for j in range(min(len(row_data), n_cols)):
                cell = tbl.cell(i + 1, j)
                para = cell.text_frame.paragraphs[0]
                para.clear()
                run = para.add_run()
                run.text = str(row_data[j])
                run.font.size = Pt(self._theme.fonts.body_size - 4)

        if elem.caption:
            self._add_placeholder(
                slide, p,   # type: ignore[arg-type]
                "",
                caption_only=elem.caption,
                caption_y_offset=p.h + 0.1,
            )

    # ── speaker note & footer ─────────────────────────────────────────────────

    def _add_speaker_note(self, slide: object, note: str) -> None:
        slide.notes_slide.notes_text_frame.text = note  # type: ignore[union-attr]

    def _add_footer(self, slide: object, override: str | None) -> None:
        footer = self._theme.footer
        text = override if override is not None else footer.custom_text
        if not text and not footer.show_page_number:
            return

        foot = slide.shapes.add_textbox(  # type: ignore[union-attr]
            Cm(2.0),
            Cm(_SLIDE_HEIGHT_CM - 1.2),
            Cm(29.87),
            Cm(0.8),
        )
        tf = foot.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(self._theme.fonts.caption_size)
        run.font.color.rgb = _rgb(self._theme.colors.text_light)

    # ── placeholder ───────────────────────────────────────────────────────────

    def _add_placeholder(
        self,
        slide: object,
        position: object,
        label: str,
        caption_only: str = "",
        caption_y_offset: float = 0.0,
    ) -> None:
        p = position
        if label:
            box = slide.shapes.add_textbox(  # type: ignore[union-attr]
                Cm(p.x), Cm(p.y), Cm(p.w), Cm(p.h)  # type: ignore[union-attr]
            )
            box.text_frame.text = label
        if caption_only:
            cap = slide.shapes.add_textbox(  # type: ignore[union-attr]
                Cm(p.x),                              # type: ignore[union-attr]
                Cm(p.y + caption_y_offset),           # type: ignore[union-attr]
                Cm(p.w),                              # type: ignore[union-attr]
                Cm(0.9),
            )
            run = cap.text_frame.paragraphs[0].add_run()
            run.text = caption_only
            run.font.size = Pt(self._theme.fonts.caption_size)
            run.font.italic = True
            run.font.color.rgb = _rgb(self._theme.colors.text_light)
