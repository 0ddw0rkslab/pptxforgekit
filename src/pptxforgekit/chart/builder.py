from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

import pandas as pd
from pptx.chart.data import ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.slide import Slide
from pptx.util import Cm, Pt

from pptxforgekit.chart.loader import ChartDataLoader
from pptxforgekit.exceptions import ChartBuildError
from pptxforgekit.models.schema import ChartElement
from pptxforgekit.models.theme import ThemeConfig

logger = logging.getLogger(__name__)

# ── PowerPoint chart type mapping ─────────────────────────────────────────────
# "bar" + vertical  → COLUMN (vertical bars)
# "bar" + horizontal → BAR (horizontal bars)
# stacked variants handled separately

_COLUMN_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED
_COLUMN_STACKED = XL_CHART_TYPE.COLUMN_STACKED
_BAR_TYPE = XL_CHART_TYPE.BAR_CLUSTERED
_BAR_STACKED = XL_CHART_TYPE.BAR_STACKED
_LINE_TYPE = XL_CHART_TYPE.LINE_MARKERS
_LINE_STACKED = XL_CHART_TYPE.LINE_MARKERS_STACKED
_SCATTER_TYPE = XL_CHART_TYPE.XY_SCATTER
_PIE_TYPE = XL_CHART_TYPE.PIE

_LEGEND_POSITION_MAP: dict[str, Any] = {
    "right":  XL_LEGEND_POSITION.RIGHT,
    "left":   XL_LEGEND_POSITION.LEFT,
    "top":    XL_LEGEND_POSITION.TOP,
    "bottom": XL_LEGEND_POSITION.BOTTOM,
}

_HAS_CATEGORY_AXIS = {"bar", "line"}   # chart types with a category + value axis pair


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class ChartBuilder:
    """Builds PowerPoint-native charts from a ChartElement schema.

    All charts are created using python-pptx's Chart API so the underlying
    data table remains editable in PowerPoint — no images are produced.
    """

    def __init__(self) -> None:
        self._loader = ChartDataLoader()

    # ── public API ────────────────────────────────────────────────────────────

    def add_chart_to_slide(
        self,
        slide: Slide,
        element: ChartElement,
        theme: ThemeConfig,
        base_path: Path | None = None,
    ) -> None:
        df = self._loader.load(element, base_path)
        x_col, y_cols = self._resolve_columns(element, df)

        chart_type = self._get_pptx_chart_type(element)
        p = element.position
        x, y, cx, cy = Cm(p.x), Cm(p.y), Cm(p.w), Cm(p.h)

        if element.chart_type == "scatter":
            chart_data = self._make_scatter_data(df, x_col, y_cols, element)
        else:
            chart_data = self._make_category_data(df, x_col, y_cols, element)

        graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
        chart = graphic_frame.chart

        self._apply_chart_title(chart, element)
        self._apply_axis(chart, element, x_col)
        self._apply_legend(chart, element, len(y_cols))
        self._apply_series_names(chart, element)
        self._apply_colors(chart, element, theme)
        self._apply_data_labels(chart, element, theme)

    # ── chart type resolution ─────────────────────────────────────────────────

    def _get_pptx_chart_type(self, element: ChartElement) -> Any:
        t = element.chart_type
        if t == "bar":
            if element.bar_direction == "horizontal":
                return _BAR_STACKED if element.stacked else _BAR_TYPE
            return _COLUMN_STACKED if element.stacked else _COLUMN_TYPE
        if t == "line":
            return _LINE_STACKED if element.stacked else _LINE_TYPE
        if t == "scatter":
            return _SCATTER_TYPE
        if t == "pie":
            return _PIE_TYPE
        raise ChartBuildError(f"Unknown chart_type '{t}'.")

    # ── column resolution ─────────────────────────────────────────────────────

    def _resolve_columns(
        self, element: ChartElement, df: pd.DataFrame
    ) -> tuple[str, list[str]]:
        all_cols = list(df.columns)

        # X column
        if element.x_column and element.x_column in all_cols:
            x_col = element.x_column
        else:
            x_col = all_cols[0]
            if element.x_column:
                logger.warning(
                    "Chart '%s': x_column '%s' not found; using '%s'.",
                    element.element_id, element.x_column, x_col,
                )

        # Y columns
        if element.y_columns and all(c in all_cols for c in element.y_columns):
            y_cols = element.y_columns
        else:
            y_cols = [
                c for c in all_cols
                if c != x_col and pd.api.types.is_numeric_dtype(df[c])
            ] or [c for c in all_cols if c != x_col]

        if not y_cols:
            raise ChartBuildError(
                f"Chart '{element.element_id}': no usable Y columns in data "
                f"(columns: {all_cols})."
            )
        return x_col, y_cols

    # ── ChartData builders ────────────────────────────────────────────────────

    def _make_category_data(
        self,
        df: pd.DataFrame,
        x_col: str,
        y_cols: list[str],
        element: ChartElement,
    ) -> ChartData:
        cd = ChartData()
        cd.categories = df[x_col].astype(str).tolist()
        display_names = element.series_names or y_cols
        for i, col in enumerate(y_cols):
            name = display_names[i] if i < len(display_names) else col
            values: list[Any] = [
                None if pd.isna(v) else float(v)
                for v in df[col].tolist()
            ]
            cd.add_series(name, tuple(values))
        return cd

    def _make_scatter_data(
        self,
        df: pd.DataFrame,
        x_col: str,
        y_cols: list[str],
        element: ChartElement,
    ) -> XyChartData:
        cd = XyChartData()
        display_names = element.series_names or y_cols
        for i, col in enumerate(y_cols):
            name = display_names[i] if i < len(display_names) else col
            series = cd.add_series(name)
            for x_val, y_val in zip(df[x_col].tolist(), df[col].tolist()):
                if pd.isna(x_val) or pd.isna(y_val):
                    continue
                series.add_data_point(float(x_val), float(y_val))
        return cd

    # ── post-creation styling ─────────────────────────────────────────────────

    def _apply_chart_title(self, chart: Any, element: ChartElement) -> None:
        if element.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = element.title

    def _apply_axis(
        self, chart: Any, element: ChartElement, x_col: str
    ) -> None:
        if element.chart_type not in _HAS_CATEGORY_AXIS:
            return

        # Category axis label
        try:
            if element.x_label:
                cat_ax = chart.category_axis
                cat_ax.has_title = True
                cat_ax.axis_title.text_frame.text = element.x_label
        except Exception as exc:
            logger.debug("Could not set category axis title: %s", exc)

        # Value axis
        try:
            val_ax = chart.value_axis

            # Label (y_label + y_unit)
            y_title = _combine_label(element.y_label, element.y_unit)
            if y_title:
                val_ax.has_title = True
                val_ax.axis_title.text_frame.text = y_title

            # Scale
            if element.value_min is not None:
                val_ax.minimum_scale = element.value_min
            if element.value_max is not None:
                val_ax.maximum_scale = element.value_max

            # Number format
            if element.number_format and element.number_format != "General":
                val_ax.tick_labels.number_format = element.number_format
                val_ax.tick_labels.number_format_is_linked = False

        except Exception as exc:
            logger.debug("Could not configure value axis: %s", exc)

    def _apply_legend(
        self, chart: Any, element: ChartElement, n_series: int
    ) -> None:
        show = n_series > 1 and element.legend_position != "none"
        chart.has_legend = show
        if show:
            pos = _LEGEND_POSITION_MAP.get(element.legend_position)
            if pos is not None:
                try:
                    chart.legend.position = pos
                    chart.legend.include_in_layout = False
                except Exception as exc:
                    logger.debug("Could not set legend position: %s", exc)

    def _apply_series_names(self, chart: Any, element: ChartElement) -> None:
        if not element.series_names:
            return
        try:
            for i, series in enumerate(chart.series):
                if i < len(element.series_names) and element.series_names[i]:
                    series.name = element.series_names[i]
        except Exception as exc:
            logger.debug("Could not rename series: %s", exc)

    def _apply_colors(
        self, chart: Any, element: ChartElement, theme: ThemeConfig
    ) -> None:
        per_series = list(element.series_colors)
        theme_colors = list(theme.chart_style.color_sequence)

        try:
            if element.chart_type == "pie":
                # Pie: one series, color each point (slice) from the theme palette
                series = chart.series[0]
                for j, point in enumerate(series.points):
                    if j < len(per_series) and per_series[j]:
                        color = per_series[j]
                    elif j < len(theme_colors):
                        color = theme_colors[j]
                    else:
                        continue
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = _rgb(color)
            else:
                for i, series in enumerate(chart.series):
                    if i < len(per_series) and per_series[i]:
                        color = per_series[i]
                    elif i < len(theme_colors):
                        color = theme_colors[i]
                    else:
                        continue
                    if element.chart_type in ("line", "scatter"):
                        # Line/scatter: color the stroke; also tint the marker fill
                        series.format.line.color.rgb = _rgb(color)
                        try:
                            series.marker.format.fill.solid()
                            series.marker.format.fill.fore_color.rgb = _rgb(color)
                        except Exception:
                            pass
                    else:
                        # Bar/column: color the series fill
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = _rgb(color)
        except Exception as exc:
            logger.debug("Could not apply series colors: %s", exc)

    def _apply_data_labels(
        self, chart: Any, element: ChartElement, theme: ThemeConfig
    ) -> None:
        show = element.show_data_labels
        if show is None:
            show = theme.chart_style.show_data_labels
        if not show:
            return

        fmt = element.data_label_format or (
            element.number_format if element.number_format != "General" else ""
        )

        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            dl = plot.data_labels
            if fmt:
                dl.number_format = fmt
                dl.number_format_is_linked = False
            dl.font.size = Pt(10)
        except Exception as exc:
            logger.debug("Could not apply data labels: %s", exc)


# ── module-level helpers ──────────────────────────────────────────────────────


def _combine_label(label: str, unit: str) -> str:
    if label and unit:
        return f"{label} {unit}"
    return label or unit
