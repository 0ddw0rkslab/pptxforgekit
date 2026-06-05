from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Cm

from pptxforgekit.chart.builder import ChartBuilder
from pptxforgekit.models.schema import ChartElement, Position
from pptxforgekit.models.theme import ThemeConfig

FIXTURES = Path(__file__).parents[1] / "fixtures" / "data"

# ── helpers ───────────────────────────────────────────────────────────────────


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _elem(chart_type: str = "bar", **kw) -> ChartElement:
    defaults = dict(
        element_id="test_chart",
        chart_type=chart_type,
        data_inline=[
            {"cat": "A", "val": 10.0, "val2": 5.0},
            {"cat": "B", "val": 20.0, "val2": 15.0},
            {"cat": "C", "val": 15.0, "val2": 10.0},
        ],
        x_column="cat",
        y_columns=["val"],
        position=Position(x=2.0, y=3.0, w=25.0, h=13.0),
    )
    defaults.update(kw)
    return ChartElement(**defaults)


def _chart_shapes(slide):
    return [s for s in slide.shapes if s.has_chart]


# ── basic rendering ───────────────────────────────────────────────────────────


class TestChartBuilderBasic:
    def test_bar_chart_creates_shape(self) -> None:
        _, slide = _blank_slide()
        ChartBuilder().add_chart_to_slide(slide, _elem("bar"), ThemeConfig())
        shapes = _chart_shapes(slide)
        assert len(shapes) == 1

    def test_chart_is_native_not_image(self) -> None:
        _, slide = _blank_slide()
        ChartBuilder().add_chart_to_slide(slide, _elem("bar"), ThemeConfig())
        chart_shape = _chart_shapes(slide)[0]
        assert chart_shape.has_chart
        # native chart has a chart attribute
        assert chart_shape.chart is not None

    def test_series_count_matches_y_columns(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", y_columns=["val", "val2"])
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        chart = _chart_shapes(slide)[0].chart
        assert len(list(chart.series)) == 2

    def test_series_data_matches_inline(self) -> None:
        _, slide = _blank_slide()
        ChartBuilder().add_chart_to_slide(slide, _elem("bar"), ThemeConfig())
        chart = _chart_shapes(slide)[0].chart
        series = list(chart.series)[0]
        # python-pptx exposes series values via .values
        values = list(series.values)
        assert values == [10.0, 20.0, 15.0]


# ── chart types ───────────────────────────────────────────────────────────────


class TestChartTypes:
    def test_line_chart(self) -> None:
        _, slide = _blank_slide()
        ChartBuilder().add_chart_to_slide(slide, _elem("line"), ThemeConfig())
        assert len(_chart_shapes(slide)) == 1

    def test_pie_chart(self) -> None:
        _, slide = _blank_slide()
        ChartBuilder().add_chart_to_slide(slide, _elem("pie"), ThemeConfig())
        assert len(_chart_shapes(slide)) == 1

    def test_scatter_chart(self) -> None:
        _, slide = _blank_slide()
        e = ChartElement(
            element_id="scatter",
            chart_type="scatter",
            data_inline=[
                {"x": 1.0, "y": 2.0},
                {"x": 3.0, "y": 5.0},
                {"x": 5.0, "y": 4.0},
            ],
            x_column="x",
            y_columns=["y"],
            position=Position(x=2.0, y=3.0, w=25.0, h=13.0),
        )
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        assert len(_chart_shapes(slide)) == 1

    def test_horizontal_bar(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", bar_direction="horizontal")
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        assert len(_chart_shapes(slide)) == 1

    def test_stacked_column(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", stacked=True, y_columns=["val", "val2"])
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        assert len(_chart_shapes(slide)) == 1


# ── titles and labels ─────────────────────────────────────────────────────────


class TestChartTitles:
    def test_chart_title_set(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", title="My Chart Title")
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        chart = _chart_shapes(slide)[0].chart
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "My Chart Title"

    def test_no_title_when_empty(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", title="")
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        chart = _chart_shapes(slide)[0].chart
        assert not chart.has_title


# ── series customisation ──────────────────────────────────────────────────────


class TestSeriesCustomisation:
    def test_series_names_override(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", y_columns=["val", "val2"], series_names=["Accuracy", "F1"])
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        chart = _chart_shapes(slide)[0].chart
        names = [s.name for s in chart.series]
        assert "Accuracy" in names
        assert "F1" in names

    def test_series_colors_applied(self) -> None:
        # Colours are applied without error (exact verification requires XML inspection)
        _, slide = _blank_slide()
        e = _elem("bar", series_colors=["#FF0000", "#00FF00"])
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        assert len(_chart_shapes(slide)) == 1


# ── legend ────────────────────────────────────────────────────────────────────


class TestLegend:
    def test_legend_shown_for_multiple_series(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", y_columns=["val", "val2"], legend_position="bottom")
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        chart = _chart_shapes(slide)[0].chart
        assert chart.has_legend

    def test_legend_hidden_when_single_series(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", y_columns=["val"])
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        chart = _chart_shapes(slide)[0].chart
        assert not chart.has_legend

    def test_legend_hidden_when_none_position(self) -> None:
        _, slide = _blank_slide()
        e = _elem("bar", y_columns=["val", "val2"], legend_position="none")
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())
        chart = _chart_shapes(slide)[0].chart
        assert not chart.has_legend


# ── file-based data ───────────────────────────────────────────────────────────


class TestFileBasedChart:
    def test_loads_csv_and_renders(self) -> None:
        _, slide = _blank_slide()
        e = ChartElement(
            element_id="csv_chart",
            chart_type="bar",
            data_source="comparison.csv",
            x_column="method",
            y_columns=["accuracy_pct"],
            title="Accuracy Comparison",
            position=Position(x=2.0, y=3.0, w=25.0, h=13.0),
        )
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig(), base_path=FIXTURES)
        chart = _chart_shapes(slide)[0].chart
        assert len(list(chart.series)) == 1

    def test_pptx_roundtrip_chart_editable(self, tmp_path: Path) -> None:
        """Saving and re-opening the PPTX still has a chart shape (not an image)."""
        prs, slide = _blank_slide()
        e = _elem("bar")
        ChartBuilder().add_chart_to_slide(slide, e, ThemeConfig())

        out = tmp_path / "test.pptx"
        prs.save(str(out))

        prs2 = Presentation(str(out))
        slide2 = prs2.slides[0]
        charts = [s for s in slide2.shapes if s.has_chart]
        assert len(charts) == 1
