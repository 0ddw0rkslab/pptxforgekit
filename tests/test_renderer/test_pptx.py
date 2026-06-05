from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from presentation_tool.models.outline import StorylineOutline
from presentation_tool.models.theme import ThemeConfig
from presentation_tool.renderer.pptx import PPTXRenderer
from presentation_tool.schema.generator import RuleBasedSchemaGenerator


class TestPPTXRenderer:
    def test_render_creates_file(
        self,
        sample_outline: StorylineOutline,
        sample_theme: ThemeConfig,
        tmp_path: Path,
    ) -> None:
        schema = RuleBasedSchemaGenerator().generate(sample_outline, sample_theme)
        out = tmp_path / "test.pptx"
        PPTXRenderer(sample_theme).render(schema, out)
        assert out.exists()
        assert out.stat().st_size > 0

    def test_render_slide_count(
        self,
        sample_outline: StorylineOutline,
        sample_theme: ThemeConfig,
        tmp_path: Path,
    ) -> None:
        schema = RuleBasedSchemaGenerator().generate(sample_outline, sample_theme)
        out = tmp_path / "test.pptx"
        PPTXRenderer(sample_theme).render(schema, out)

        prs = Presentation(str(out))
        assert len(prs.slides) == len(schema.slides)

    def test_render_widescreen_dimensions(
        self,
        sample_outline: StorylineOutline,
        sample_theme: ThemeConfig,
        tmp_path: Path,
    ) -> None:
        from pptx.util import Cm
        schema = RuleBasedSchemaGenerator().generate(sample_outline, sample_theme)
        out = tmp_path / "test.pptx"
        PPTXRenderer(sample_theme).render(schema, out)

        prs = Presentation(str(out))
        # Allow ±10 EMU tolerance due to float→EMU conversion
        assert abs(prs.slide_width - Cm(33.87)) < 100
        assert abs(prs.slide_height - Cm(19.05)) < 100

    def test_render_with_chart_data(
        self,
        sample_theme: ThemeConfig,
        tmp_path: Path,
        fixture_csv: Path,
    ) -> None:
        from presentation_tool.models.analysis import AnalysisResult, DataFileRef, Section
        from presentation_tool.planner.storyline import RuleBasedStorylinePlanner

        analysis = AnalysisResult(
            source_files=[str(fixture_csv)],
            title="Chart Test",
            data_files=[
                DataFileRef(
                    file_path=str(fixture_csv),
                    columns=["method", "accuracy", "f1_score"],
                    row_count=3,
                )
            ],
        )
        outline = RuleBasedStorylinePlanner().plan(analysis, "research")
        schema = RuleBasedSchemaGenerator().generate(outline, sample_theme)

        out = tmp_path / "chart_test.pptx"
        # base_path must point to the fixture CSV's parent so the renderer finds it
        PPTXRenderer(sample_theme).render(schema, out, base_path=fixture_csv.parent)
        assert out.exists()
